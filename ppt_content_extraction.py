import os
from pptx import Presentation

def extract_text(ppt_path):
    """Extract text from a PowerPoint presentation"""
    prs = Presentation(ppt_path)
    slides_text = {}

    for i, slide in enumerate(prs.slides, start=1):
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
        slides_text[f"Slide {i}"] = text

    return slides_text

def extract_images(ppt_path, output_folder):
    """Extract images from a PowerPoint presentation and save them"""
    prs = Presentation(ppt_path)
    os.makedirs(output_folder, exist_ok=True)
    
    image_list = []
    image_count = 0

    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 corresponds to a picture
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext  # Get image extension
                
                image_filename = f"{output_folder}/slide_{i}_image_{image_count + 1}.{image_ext}"
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes)
                
                image_list.append(image_filename)
                image_count += 1
    
    return image_list

def extract_tables(ppt_path):
    """Extract tables with metadata from a PowerPoint presentation"""
    prs = Presentation(ppt_path)
    tables_metadata = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                num_rows = len(table.rows)
                num_columns = len(table.columns)
                table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]

                tables_metadata.append({
                    "slide_number": slide_num,
                    "num_rows": num_rows,
                    "num_columns": num_columns,
                    "data": table_data
                })
    
    return tables_metadata

def extract_links(ppt_path):
    """Extract hyperlinks from a PowerPoint presentation"""
    prs = Presentation(ppt_path)
    links = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink and run.hyperlink.address:
                            links.append({"slide_number": slide_num, "link": run.hyperlink.address})
    
    return links

# Example usage
ppt_path = "ppt_test.pptx"  # Replace with your PowerPoint file
output_folder = "ppt_images"

# Extract text
ppt_text = extract_text(ppt_path)
print("Extracted Text:")
for slide, text in ppt_text.items():
    print(f"{slide}:\n{text}\n")

# Extract images
images = extract_images(ppt_path, output_folder)
print("Extracted Images:", images)

# Extract tables
tables = extract_tables(ppt_path)
print("Extracted Tables : ")
for table in tables:
    print(f"Table on Slide {table['slide_number']}: Rows={table['num_rows']}, Columns={table['num_columns']}")
    for row in table["data"]:
        print(row)
    print()


# Extract links
links = extract_links(ppt_path)
print("Extracted Links:")
for link in links:
    print(f"Slide {link['slide_number']} -> {link['link']}")

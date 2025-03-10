from abc import ABC, abstractmethod
import os
import docx.document
import fitz  # PyMuPDF for PDF processing
import docx
import pptx
import mysql.connector
import csv
import pdfplumber
from io import BytesIO
import pandas as pd
from docx.opc.constants import RELATIONSHIP_TYPE
import pptx.presentation
from pptx import Presentation
import json
from PIL import Image





# Abstract class for file loading
class FileLoader(ABC):
    def __init__(self, file_path : str):
        self.file_path = file_path
        self.validate_file()

    @abstractmethod
    def validate_file(self):
        pass

    @abstractmethod
    def load_file(self):
        pass

# Concrete class for PDF file loading
class PDFLoader(FileLoader):
    def validate_file(self):
        if not self.file_path.lower().endswith(".pdf"):
            raise ValueError("Invalid PDF file format")
    
    def load_file(self):
        return fitz.open(self.file_path)

# Concrete class for DOCX file loading
class DOCXLoader(FileLoader):
    def validate_file(self):
        if not self.file_path.lower().endswith(".docx"):
            raise ValueError("Invalid DOCX file format")
    
    def load_file(self):
        return docx.Document(self.file_path)

# Concrete class for PPT file loading
class PPTLoader(FileLoader):
    def validate_file(self):
        if not self.file_path.lower().endswith(".pptx"):
            raise ValueError("Invalid PPT file format")
    
    def load_file(self):
        return pptx.Presentation(self.file_path)

# Data extractor class
class DataExtractor:
    def __init__(self, file_loader : FileLoader):
        self.file_loader = file_loader.load_file()
        self.file_path = file_loader.file_path


    def extract_text(self):
        extracted_text = []
        
        if isinstance(self.file_loader, fitz.Document):  # PDF
            for page_num, page in enumerate(self.file_loader, start=1):
                page_data = {"page_number": page_num, "text": [], "metadata": []}
                
                text_blocks = page.get_text("dict")["blocks"]  # Extract text blocks
                
                for block in text_blocks:
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:  # Spans contain individual text segments
                                text_content = span["text"].strip()
                                if text_content:
                                    page_data["text"].append(text_content)
                                    page_data["metadata"].append({
                                        "text": text_content,
                                        "font_size": span["size"],
                                        "bold": "Bold" in span["font"],
                                        "italic": "Italic" in span["font"]
                                    })
                
                extracted_text.append(page_data)

        elif isinstance(self.file_loader, docx.document.Document):  # DOCX
            page_data = {"text": [], "metadata": []}
            
            for para in self.file_loader.paragraphs:
                text_content = para.text.strip()
                if text_content:
                    font_size = None
                    is_bold = False
                    is_italic = False
                    
                    if para.runs:
                        font_size = para.runs[0].font.size.pt if para.runs[0].font.size else None
                        is_bold = para.runs[0].bold
                        is_italic = para.runs[0].italic
                    
                    is_heading = para.style.name.startswith("Heading")
                    
                    page_data["text"].append(text_content)
                    page_data["metadata"].append({
                        "text": text_content,
                        "font_size": font_size,
                        "bold": is_bold,
                        "italic": is_italic,
                        "heading": is_heading
                    })
            
            extracted_text.append(page_data)

        elif isinstance(self.file_loader, pptx.presentation.Presentation):  # PPTX
            for slide_num, slide in enumerate(self.file_loader.slides, start=1):
                slide_data = {"slide_number": slide_num, "text": [], "metadata": []}
                
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for para in shape.text_frame.paragraphs:
                            text_content = para.text.strip()
                            if text_content:
                                font_size = para.runs[0].font.size.pt if para.runs and para.runs[0].font.size else None
                                
                                slide_data["text"].append(text_content)
                                slide_data["metadata"].append({
                                    "text": text_content,
                                    "font_size": font_size,
                                })
                
                extracted_text.append(slide_data)
        
        return extracted_text



    


    def extract_links(self):
        links = []
        if isinstance(self.file_loader, fitz.Document):  # PDF
            for page_num, page in enumerate(self.file_loader):
                for link in page.get_links():
                    links.append((page_num + 1, link.get("uri", "")))
        elif isinstance(self.file_loader, docx.document.Document):  # DOCX
            for rel in self.file_loader.part.rels:
                if "hyperlink" in self.file_loader.part.rels[rel].reltype:
                    links.append(self.file_loader.part.rels[rel].target_ref)
        elif isinstance(self.file_loader,pptx.presentation.Presentation):
            prs = self.file_loader
            links = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    links.append({"slide_number": slide_num, "link": run.hyperlink.address})
            
            return links
        return links

    

    def extract_images(self):
        images_with_metadata = []

        if isinstance(self.file_loader, fitz.Document):  # PDF
            doc = self.file_loader
            output_folder = "pdf_images"
            for page_num, page in enumerate(doc, start=1):
                for img_index, img in enumerate(page.get_images(full=True), start=1):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]

                    image_filename = f"{output_folder}/pdf_page_{page_num}_image_{img_index}.{image_ext}"
                    with open(image_filename, "wb") as img_file:
                        img_file.write(image_bytes)

                    # Get resolution
                    with Image.open(image_filename) as img:
                        width, height = img.size

                    images_with_metadata.append({
                        "filename": image_filename,
                        "page_number": page_num,
                        "width": width,
                        "height": height
                    })

        elif isinstance(self.file_loader, docx.document.Document):  # DOCX
            doc = self.file_loader
            output_folder = "docx_images"
            for rel_id, rel in doc.part.rels.items():
                if rel.reltype == RELATIONSHIP_TYPE.IMAGE:
                    image_part = rel.target_part
                    image_data = image_part.blob
                    image_ext = image_part.content_type.split("/")[-1]  # Get extension

                    image_filename = f"{output_folder}/image_{len(images_with_metadata) + 1}.{image_ext}"
                    with open(image_filename, "wb") as img_file:
                        img_file.write(image_data)

                    # Get resolution
                    with Image.open(image_filename) as img:
                        width, height = img.size

                    images_with_metadata.append({
                        "filename": image_filename,
                        "width": width,
                        "height": height
                    })

        elif isinstance(self.file_loader, pptx.presentation.Presentation):  # PPTX
            output_folder = "ppt_images"
            for slide_num, slide in enumerate(self.file_loader.slides, start=1):
                image_count = 1
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Shape type 13 = Picture
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext  # Get image extension

                        image_filename = f"{output_folder}/ppt_page_{slide_num}_image_{image_count}.{image_ext}"
                        with open(image_filename, "wb") as img_file:
                            img_file.write(image_bytes)

                        # Get resolution
                        with Image.open(image_filename) as img:
                            width, height = img.size

                        images_with_metadata.append({
                            "filename": image_filename,
                            "page_number": slide_num,  # Treating slides as "pages"
                            "width": width,
                            "height": height
                        })

                        image_count += 1

        return images_with_metadata


    def extract_tables(self):
        tables_with_metadata = []

        if isinstance(self.file_loader, docx.document.Document):  # DOCX
            for table in self.file_loader.tables:
                data = [[cell.text for cell in row.cells] for row in table.rows]
                tables_with_metadata.append(pd.DataFrame(data))

        elif isinstance(self.file_loader, fitz.Document):  # PDF
            with pdfplumber.open(self.file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    extracted_table = page.extract_table()
                    if extracted_table:
                        df = pd.DataFrame(extracted_table[1:], columns=extracted_table[0])
                        df.insert(0, "Page Number", page_num)  # Add page number as a column
                        tables_with_metadata.append(df)

        elif isinstance(self.file_loader, pptx.presentation.Presentation):  # PPTX
            for slide_num, slide in enumerate(self.file_loader.slides, start=1):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                        df = pd.DataFrame(data)
                        df.insert(0, "Slide Number", slide_num)  # Add slide number as a column
                        tables_with_metadata.append(df)

        return tables_with_metadata


# Abstract storage class
class Storage(ABC):
    @abstractmethod
    def store(self, data):
        pass

# Concrete class for file storage
class FileStorage(Storage):
    def store(self, data):
        if isinstance(data, str):
            with open("extracted_text.txt", "w", encoding="utf-8") as f:
                f.write(data)
        elif isinstance(data, list):
            for i, item in enumerate(data):
                if isinstance(item, pd.DataFrame):
                    item.to_csv(f"extracted_table_{i}.csv", index=False)
                else:
                    with open(f"extracted_data_{i}.txt", "w", encoding="utf-8") as f:
                        f.write(str(item))



class SQLStorage(Storage):
    def __init__(self, host="localhost", user="root", password="rootroot", database="assignment3"):
        self.conn = mysql.connector.connect(
            host=host, user=user, password=password, database=database
        )
        self.cursor = self.conn.cursor()
        self.create_table()

    def create_table(self):
        self.cursor.execute("""
        CREATE TABLE IF NOT EXISTS extracted_data (
            id INT AUTO_INCREMENT PRIMARY KEY,
            file_name VARCHAR(255),
            text_content TEXT,
            links JSON,
            tables JSON,
            images JSON
        )
        """)
        self.conn.commit()

    def store(self, file_name, text, links, tables, images):
        # Convert lists/dicts to JSON for storage
        links_json = json.dumps(links, default=str) if links else "[]"
        tables_json = json.dumps([df.to_dict(orient="records") for df in tables], default=str) if tables else "[]"
        images_json = json.dumps(images, default=str) if images else "[]"

        text_str = json.dumps(text, default=str)  # Convert list to JSON string

        self.cursor.execute("""
        INSERT INTO extracted_data (file_name, text_content, links, tables, images)
        VALUES (%s, %s, %s, %s, %s)
        """, (file_name, text_str, links_json, tables_json, images_json))
        
        self.conn.commit()

# Main function for testing
def main():
    # file_path = "sample_pdfs/test1.pdf"
    file_path = "sample_docx/demo.docx"  # Change this to test different file types
    # file_path = "sample_pptx/ppt_test.pptx"


    if file_path.endswith(".pdf"):
        loader = PDFLoader(file_path)
    elif file_path.endswith(".docx"):
        loader = DOCXLoader(file_path)
    elif file_path.endswith(".pptx"):
        loader = PPTLoader(file_path)
    else:
        print("Unsupported file format.")
        return
    
    extractor = DataExtractor(loader)
    text = extractor.extract_text()
    links = extractor.extract_links()
    images = extractor.extract_images()
    tables = extractor.extract_tables()
    
    print("Extracted Text:\n", text)
    print("Extracted Links:\n", links)
    print("Extracted Images:", images, "images found")
    print("Extracted Tables:", tables, "tables found")
    
    
    file_storage = FileStorage()
    mysql_storage = SQLStorage()
    
    # Store in File Storage
    file_storage.store(text)
    file_storage.store(tables)

    mysql_storage.store(os.path.basename(file_path), text, links, tables, images)
    
    print("Data stored successfully.")

if __name__ == "__main__":
    main()




























































# Concrete class for MySQL storage
# class SQLStorage(Storage):
#     def __init__(self, host="localhost", user="root", password="rootroot", database="assignment3"):
#         self.conn = mysql.connector.connect(
#             host=host, user=user, password=password, database=database
#         )
#         self.cursor = self.conn.cursor()
#         self.create_tables()

#     def create_tables(self):
#         self.cursor.execute("""
#         CREATE TABLE IF NOT EXISTS extracted_text (
#             id INT AUTO_INCREMENT PRIMARY KEY,
#             content TEXT
#         )
#         """)
#         self.cursor.execute("""
#         CREATE TABLE IF NOT EXISTS extracted_tables (
#             id INT AUTO_INCREMENT PRIMARY KEY,
#             source VARCHAR(10),
#             table_data JSON
#         )
#         """)
#         self.conn.commit()

#     def store(self, data):
#         if isinstance(data, str):  # Store extracted text
#             self.cursor.execute("INSERT INTO extracted_text (content) VALUES (%s)", (data,))
        
#         elif isinstance(data, list):  # Store extracted tables
#             for item in data:
#                 if isinstance(item, pd.DataFrame):  # Store DataFrame from PDF/DOCX
#                     table_json = item.to_json(orient="records")  # Convert DataFrame to JSON
#                     self.cursor.execute("INSERT INTO extracted_tables (source, table_data) VALUES (%s, %s)", ("docx", table_json))
#                 elif isinstance(item, dict):  # Store PPT tables
#                     table_json = json.dumps(item, default=str)  # Convert dictionary to JSON
#                     self.cursor.execute("INSERT INTO extracted_tables (source, table_data) VALUES (%s, %s)", ("ppt/pdf", table_json))
    
#         self.conn.commit()
































    # def extract_text(self):
    #     extracted_text = ""

    #     if isinstance(self.file_loader, fitz.Document):  # PDF
    #         for page_num, page in enumerate(self.file_loader, start=1):
    #             extracted_text += f"\n--- Page {page_num} ---\n"
    #             extracted_text += page.get_text("text") + "\n"

    #     elif isinstance(self.file_loader, docx.document.Document):  # DOCX
    #         extracted_text += "\n--- DOCX Content ---\n"
    #         for i, para in enumerate(self.file_loader.paragraphs, start=1):
    #             extracted_text += f"\n--- Paragraph {i} ---\n{para.text}\n"

    #     elif isinstance(self.file_loader, pptx.presentation.Presentation):  # PPTX
    #         for slide_num, slide in enumerate(self.file_loader.slides, start=1):
    #             extracted_text += f"\n--- Slide {slide_num} ---\n"
    #             for shape in slide.shapes:
    #                 if hasattr(shape, "text") and shape.text.strip():
    #                     extracted_text += shape.text + "\n"

    #     return extracted_text